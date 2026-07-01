"""PowerPoint 보고서 자동 생성."""
from __future__ import annotations

from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt

from quality_mh.ai_analyzer import analyze_top5_mh
from quality_mh.models import CalcResult, QualitativeRecord
from quality_mh.plant_config import PlantConfig
from quality_mh.summary_engine import SummaryReport, build_summary_report


def _add_title_slide(prs: Presentation, config: PlantConfig) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{config.plant_name} 품질 M/H 분석 보고서"
    slide.placeholders[1].text = (
        f"{config.analysis_year}년 | 근무시간 {config.work_hours_per_day}hr\n"
        f"생성일: {datetime.now().strftime('%Y-%m-%d')}"
    )


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    cols, row_count = len(headers), len(rows) + 1
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * row_count)
    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)


def export_ppt(
    config: PlantConfig,
    calc_results: list[CalcResult],
    qualitative: list[QualitativeRecord],
    report: SummaryReport | None = None,
) -> BytesIO:
    """분석 결과 PowerPoint 생성."""
    report = report or build_summary_report(config, calc_results, qualitative)
    prs = Presentation()
    _add_title_slide(prs, config)

    org_bullets = [
        f"공장: {config.plant_name}",
        f"분석년도: {config.analysis_year}",
        f"1인 근무시간: {config.work_hours_per_day}hr/일",
        f"연간 가용시간: {config.work_hours_per_year:,.1f}hr",
        f"부가공수: {config.allowance_rate * 100:.0f}%",
        f"교대근무: {config.shift_type}",
    ]
    _add_bullet_slide(prs, "1. 조직현황", org_bullets)

    std_rows = []
    for row in report.quantitative_rows:
        if row.sub_label in ("입고", "공정", "완성", "시험", "合"):
            std_rows.append([row.sub_label, f"{row.current:.0f}", f"{row.standard:.1f}", f"{row.diff:+.1f}"])
    if report.qualitative_row:
        q = report.qualitative_row
        std_rows.append(["정성", f"{q.current:.0f}", f"{q.standard:.0f}", f"{q.diff:+.0f}"])
    _add_table_slide(prs, "2. 표준인원 vs 현재원", ["구분", "현재원", "표준인원", "차이"], std_rows)

    gap_bullets = report.gap_comments or ["Gap이 허용 범위 내입니다."]
    _add_bullet_slide(prs, "3. Gap 분석 의견", gap_bullets)

    top5 = analyze_top5_mh(calc_results)
    top_rows = [[t["순위"], t["업무"][:30], t["표준공수"], t["표준인원"]] for t in top5]
    _add_table_slide(prs, "4. 주요업무 TOP5", ["순위", "업무", "M/H", "인원"], top_rows)

    qual_bullets = [
        f"{r.task_name}: {r.standard_headcount}명 — {r.task_definition or ''}"
        for r in qualitative[:6]
    ] or ["정성 레코드 없음"]
    _add_bullet_slide(prs, "5. 정성 업무 현황", qual_bullets)

    # 분석 의견
    analysis = [
        f"정량 업무 {len(calc_results)}건 분석 완료",
        f"총 표준공수 합계: {sum(r.standard_mh for r in calc_results):.2f} M/H",
        f"총 표준인원(정량): {sum(r.standard_headcount for r in calc_results)}명",
        f"정성 표준인원: {sum(r.standard_headcount for r in qualitative)}명",
    ]
    _add_bullet_slide(prs, "6. 종합 분석 의견", analysis)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
